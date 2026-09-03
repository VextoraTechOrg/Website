"""
Generate VextoraTech company profile PPTX from docs/company-profile-pptx-brief.md
Run: python scripts/generate-company-profile-pptx.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "VextoraTech-Company-Profile.pptx"

# Brand palette
BG = RGBColor(0x0D, 0x12, 0x19)
SURFACE = RGBColor(0x16, 0x1A, 0x20)
PRIMARY = RGBColor(0x5B, 0x9F, 0xD4)
PRIMARY_DARK = RGBColor(0x4A, 0x8F, 0xC4)
TEXT = RGBColor(0xE8, 0xEA, 0xED)
MUTED = RGBColor(0x8B, 0x94, 0x9E)
BORDER = RGBColor(0x2A, 0x31, 0x3A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.7)
MARGIN_Y = Inches(0.55)


def set_run(run, size=18, bold=False, color=TEXT, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_para(tf, text, *, size=16, bold=False, color=TEXT, space_before=6, space_after=4, font="Calibri"):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return p


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, fill=SURFACE, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_shape(shape, fill)
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_accent_rule(slide, left, top, width=Inches(0.9)):
    return add_rect(slide, left, top, width, Pt(3), fill=PRIMARY)


def add_footer(slide, page: int, total: int = 14):
    add_textbox(
        slide,
        MARGIN_X,
        Inches(7.05),
        Inches(10),
        Inches(0.3),
        "VextoraTech  ·  Lahore  ·  vextoratech.com  ·  info@vextoratech.com",
        size=10,
        color=MUTED,
    )
    add_textbox(
        slide,
        Inches(11.6),
        Inches(7.05),
        Inches(1.2),
        Inches(0.3),
        f"{page:02d} / {total:02d}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill_shape(bg, BG)
    # send background to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return slide


def section_title(slide, eyebrow, title):
    add_textbox(slide, MARGIN_X, MARGIN_Y, Inches(10), Inches(0.3), eyebrow.upper(), size=11, bold=True, color=PRIMARY)
    add_accent_rule(slide, MARGIN_X, Inches(0.9))
    add_textbox(slide, MARGIN_X, Inches(1.05), Inches(11.5), Inches(0.7), title, size=32, bold=True, color=TEXT, font="Calibri")


def card(slide, left, top, width, height, title, body, *, tag=None):
    add_rect(slide, left, top, width, height, fill=SURFACE, line=BORDER)
    y = top + Inches(0.22)
    if tag:
        add_textbox(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.25), tag.upper(), size=10, bold=True, color=PRIMARY)
        y += Inches(0.28)
    add_textbox(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.4), title, size=16, bold=True, color=TEXT)
    add_textbox(slide, left + Inches(0.25), y + Inches(0.4), width - Inches(0.5), height - (y - top) - Inches(0.55), body, size=12, color=MUTED)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 14

    # ── 01 Title ─────────────────────────────────────────────
    s = blank_slide(prs)
    add_accent_rule(s, MARGIN_X, Inches(2.1), Inches(1.2))
    add_textbox(s, MARGIN_X, Inches(2.3), Inches(12), Inches(0.9), "VextoraTech", size=48, bold=True, color=PRIMARY, font="Calibri")
    add_textbox(s, MARGIN_X, Inches(3.2), Inches(11), Inches(0.6), "We build software that thinks.", size=28, bold=True, color=TEXT)
    add_textbox(
        s,
        MARGIN_X,
        Inches(4.0),
        Inches(10),
        Inches(0.5),
        "AI · Web · Mobile · Cloud  —  engineering for ambitious teams",
        size=16,
        color=MUTED,
    )
    add_textbox(s, MARGIN_X, Inches(5.2), Inches(10), Inches(0.4), "Lahore, Pakistan  ·  Incorporated 2026  ·  vextoratech.com", size=14, color=MUTED)
    add_footer(s, 1, total)

    # ── 02 Who we are ────────────────────────────────────────
    s = blank_slide(prs)
    section_title(s, "Company", "Who we are")
    bullets = [
        "Engineering-led software studio that ships intelligent digital products — AI models and RAG systems to full-stack web, mobile, and cloud.",
        "Incorporated in 2026 in Lahore. No legacy agency, no rebrand, no borrowed history.",
        "Small team building in public from day one.",
        "Clients work directly with the engineers shipping the product — not a sales layer in between.",
        "Built by engineers, for founders.",
    ]
    box = s.shapes.add_textbox(MARGIN_X, Inches(2.0), Inches(11.5), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(10)
        p.space_after = Pt(6)
        p.level = 0
        run = p.add_run()
        run.text = b
        set_run(run, size=16, color=TEXT)
    add_footer(s, 2, total)

    # ── 03 Mission & vision ──────────────────────────────────
    s = blank_slide(prs)
    section_title(s, "Purpose", "Mission & vision")
    card(
        s,
        MARGIN_X,
        Inches(2.0),
        Inches(5.7),
        Inches(2.4),
        "Mission",
        "Close the gap between great ideas and trusted technical execution. Ship production software that is accurate, scalable, and honest about cost and risk.",
    )
    card(
        s,
        Inches(6.9),
        Inches(2.0),
        Inches(5.7),
        Inches(2.4),
        "Vision",
        "Become the technical partner ambitious teams call when they need AI, web, and infrastructure delivered as one coherent product — not as disconnected vendors.",
    )
    add_rect(s, MARGIN_X, Inches(4.7), Inches(11.9), Inches(1.5), fill=SURFACE, line=BORDER)
    add_textbox(
        s,
        MARGIN_X + Inches(0.4),
        Inches(4.95),
        Inches(11.1),
        Inches(1.1),
        '"Technology should accelerate your vision, not complicate it. We build the infrastructure that makes ambitious ideas inevitable."',
        size=16,
        color=TEXT,
    )
    add_footer(s, 3, total)

    # ── 04 Values ────────────────────────────────────────────
    s = blank_slide(prs)
    section_title(s, "Culture", "What we stand for")
    values = [
        ("Radical Honesty", "Say what you'd tell a co-founder. Clear beats comfortable."),
        ("Ownership Mentality", "Treat the product like it's ours — edge cases and follow-through included."),
        ("Technical Depth", "Pattern design, performance, security — not just \"it works.\""),
        ("Long-Term Thinking", "Decisions today should not create debt for the next 10,000 users."),
    ]
    for i, (title, body) in enumerate(values):
        col = i % 2
        row = i // 2
        left = MARGIN_X + col * Inches(6.1)
        top = Inches(2.0) + row * Inches(2.15)
        card(s, left, top, Inches(5.85), Inches(1.95), title, body)
    add_footer(s, 4, total)

    # ── 05 Domain strengths ──────────────────────────────────
    s = blank_slide(prs)
    section_title(s, "Expertise", "Where we go deep")
    domains = [
        ("Computer Vision", "Object detection, OCR, ANPR, face recognition, multi-camera surveillance — built for real-world accuracy and latency."),
        ("Voice AI", "Whisper transcription, speaker diarization, and call analytics that turn audio into actionable insight."),
        ("RAG & LLM Systems", "Grounded assistants with vector search, citations, and guardrails — factual and safe to ship."),
        ("Full-Stack Products", "React and FastAPI platforms, APIs, auth, and cloud — from prototype through launch and scale."),
    ]
    for i, (title, body) in enumerate(domains):
        top = Inches(1.95) + i * Inches(1.15)
        add_rect(s, MARGIN_X, top, Inches(11.9), Inches(1.05), fill=SURFACE, line=BORDER)
        add_rect(s, MARGIN_X, top, Pt(4), Inches(1.05), fill=PRIMARY)
        add_textbox(s, MARGIN_X + Inches(0.35), top + Inches(0.18), Inches(3.5), Inches(0.35), title, size=16, bold=True, color=TEXT)
        add_textbox(s, MARGIN_X + Inches(4.0), top + Inches(0.22), Inches(7.5), Inches(0.7), body, size=13, color=MUTED)
    add_footer(s, 5, total)

    # ── 06 Services ──────────────────────────────────────────
    s = blank_slide(prs)
    section_title(s, "Capabilities", "Services")
    services = [
        ("AI & Machine Learning", "Vision, voice, RAG, local inference, production model APIs."),
        ("Web Development", "React, Next.js, TanStack, FastAPI, Node, Postgres, RBAC."),
        ("Mobile Apps", "React Native / Flutter, offline-first, store deployment."),
        ("Cloud & DevOps", "AWS/GCP/Azure, Docker, K8s, CI/CD, Terraform, monitoring."),
        ("UI/UX Design", "Research, Figma prototypes, design systems, WCAG 2.1 AA."),
        ("API & Integrations", "OpenAPI REST, GraphQL, webhooks, third-party systems."),
    ]
    for i, (title, body) in enumerate(services):
        col = i % 3
        row = i // 3
        left = MARGIN_X + col * Inches(4.05)
        top = Inches(2.0) + row * Inches(2.2)
        card(s, left, top, Inches(3.85), Inches(2.0), title, body)
    add_footer(s, 6, total)

    # ── 07 How we work ───────────────────────────────────────
    s = blank_slide(prs)
    section_title(s, "Process", "How we work")
    steps = [
        ("01", "Discovery", "Goals, constraints, scope, and success metrics."),
        ("02", "Design", "Wireframes and Figma prototypes before heavy code."),
        ("03", "Build", "Agile sprints, weekly demos, real deploys early."),
        ("04", "Launch", "Production rollout, monitoring, and ongoing support."),
    ]
    for i, (num, title, body) in enumerate(steps):
        left = MARGIN_X + i * Inches(3.05)
        add_rect(s, left, Inches(2.3), Inches(2.9), Inches(3.4), fill=SURFACE, line=BORDER)
        add_textbox(s, left + Inches(0.25), Inches(2.55), Inches(2.4), Inches(0.4), num, size=22, bold=True, color=PRIMARY)
        add_textbox(s, left + Inches(0.25), Inches(3.2), Inches(2.4), Inches(0.4), title, size=18, bold=True, color=TEXT)
        add_textbox(s, left + Inches(0.25), Inches(3.8), Inches(2.4), Inches(1.5), body, size=13, color=MUTED)
    add_footer(s, 7, total)

    # ── 08 Engagement models ─────────────────────────────────
    s = blank_slide(prs)
    section_title(s, "Engagement", "How to work with us")
    models = [
        ("Fixed-scope build", "Defined MVP or feature set with agreed deliverables, timeline, and acceptance criteria. Best when requirements are clear."),
        ("Ongoing partnership", "Retained engineering team embedded in your roadmap — sprints, demos, continuous delivery. Best for products in active growth."),
        ("Discovery / PoC", "Short spike to validate feasibility, architecture, and cost before a full build. Low-risk way to start."),
    ]
    for i, (title, body) in enumerate(models):
        top = Inches(2.0) + i * Inches(1.45)
        add_rect(s, MARGIN_X, top, Inches(11.9), Inches(1.3), fill=SURFACE, line=BORDER)
        add_textbox(s, MARGIN_X + Inches(0.4), top + Inches(0.25), Inches(11), Inches(0.35), title, size=18, bold=True, color=PRIMARY)
        add_textbox(s, MARGIN_X + Inches(0.4), top + Inches(0.65), Inches(11), Inches(0.5), body, size=14, color=MUTED)
    add_footer(s, 8, total)

    # ── 09 Selected work ─────────────────────────────────────
    s = blank_slide(prs)
    section_title(s, "Portfolio", "Selected work")
    add_textbox(
        s,
        MARGIN_X,
        Inches(1.85),
        Inches(11.9),
        Inches(0.55),
        "Some of this was delivered by our people before VextoraTech was incorporated; some we are building now. Same engineers, same standards — under one name from 2026 onward.",
        size=12,
        color=MUTED,
    )
    projects = [
        ("NexaWatch", "Prototype", "Real-time field visibility — live GPS, geofences, route history."),
        ("NexaDesk AI", "Prototype", "Secure AI helpdesk — ticket analysis, smart replies, priority detection."),
        ("Voice Intelligence Hub", "Prior team work", "Transcribe, diarize, analyze calls — sentiment, topics, actions."),
        ("AI Surveillance System", "Prior team work", "YOLOv8 multi-camera detection — zone breaches, loitering, alerts."),
        ("PYLI", "Prior team work", "Centralized multi-profile business platform with bespoke UI."),
        ("Medical Knowledge Assistant", "Internal build", "Private RAG for clinical and document Q&A."),
    ]
    for i, (name, origin, summary) in enumerate(projects):
        col = i % 3
        row = i // 3
        left = MARGIN_X + col * Inches(4.05)
        top = Inches(2.55) + row * Inches(2.0)
        card(s, left, top, Inches(3.85), Inches(1.85), name, summary, tag=origin)
    add_footer(s, 9, total)

    # ── 10 Featured: NexaWatch ───────────────────────────────
    s = blank_slide(prs)
    section_title(s, "Featured work · Prototype", "NexaWatch")
    add_textbox(s, MARGIN_X, Inches(1.9), Inches(11.5), Inches(0.4), "Real-time field visibility for teams, routes, and client visits.", size=16, color=MUTED)
    features = [
        ("Challenge", "Field teams need live location, geofence alerts, and route history in one secure admin surface — not scattered tools."),
        ("Approach", "Built a real-time GPS tracking dashboard with geofences, route history, and role-based admin access."),
        ("Stack", "React · Node.js · WebSocket · GPS tracking"),
    ]
    for i, (title, body) in enumerate(features):
        top = Inches(2.5) + i * Inches(1.25)
        add_rect(s, MARGIN_X, top, Inches(11.9), Inches(1.1), fill=SURFACE, line=BORDER)
        add_textbox(s, MARGIN_X + Inches(0.4), top + Inches(0.2), Inches(2.2), Inches(0.35), title, size=14, bold=True, color=PRIMARY)
        add_textbox(s, MARGIN_X + Inches(2.8), top + Inches(0.22), Inches(8.6), Inches(0.7), body, size=14, color=TEXT)
    add_footer(s, 10, total)

    # ── 11 Featured: Voice Intelligence Hub ──────────────────
    s = blank_slide(prs)
    section_title(s, "Featured work · Prior team work", "Voice Intelligence Hub")
    add_textbox(s, MARGIN_X, Inches(1.9), Inches(11.5), Inches(0.4), "End-to-end voice analytics for contact centers and sales teams.", size=16, color=MUTED)
    features = [
        ("Challenge", "Raw call recordings stay unused — teams need transcription, speaker separation, and actionable insights."),
        ("Approach", "Whisper + NLP pipelines that diarize speakers, extract sentiment and topics, and surface action items."),
        ("Stack", "Whisper · Python · NLP · FastAPI"),
    ]
    for i, (title, body) in enumerate(features):
        top = Inches(2.5) + i * Inches(1.25)
        add_rect(s, MARGIN_X, top, Inches(11.9), Inches(1.1), fill=SURFACE, line=BORDER)
        add_textbox(s, MARGIN_X + Inches(0.4), top + Inches(0.2), Inches(2.2), Inches(0.35), title, size=14, bold=True, color=PRIMARY)
        add_textbox(s, MARGIN_X + Inches(2.8), top + Inches(0.22), Inches(8.6), Inches(0.7), body, size=14, color=TEXT)
    add_footer(s, 11, total)

    # ── 12 Technology ────────────────────────────────────────
    s = blank_slide(prs)
    section_title(s, "Stack", "Technology we ship with")
    stacks = [
        ("Frontend", "React, Next.js, TanStack, Tailwind CSS, Figma"),
        ("Backend", "FastAPI, Python, Node.js"),
        ("Data", "PostgreSQL, MongoDB"),
        ("AI / ML", "LLMs, Whisper, YOLO, OpenCV, PyTorch, Ollama"),
        ("Infra", "Docker, AWS, CI/CD"),
    ]
    for i, (area, techs) in enumerate(stacks):
        top = Inches(2.0) + i * Inches(0.85)
        add_rect(s, MARGIN_X, top, Inches(11.9), Inches(0.75), fill=SURFACE, line=BORDER)
        add_textbox(s, MARGIN_X + Inches(0.4), top + Inches(0.2), Inches(2.5), Inches(0.4), area, size=15, bold=True, color=PRIMARY)
        add_textbox(s, MARGIN_X + Inches(3.2), top + Inches(0.2), Inches(8.2), Inches(0.4), techs, size=15, color=TEXT)
    add_footer(s, 12, total)

    # ── 13 Team ──────────────────────────────────────────────
    s = blank_slide(prs)
    section_title(s, "People", "The team")
    add_textbox(s, MARGIN_X, Inches(1.85), Inches(11.5), Inches(0.35), "Lahore-based · remote-friendly across Pakistan · direct engineer–client collaboration", size=12, color=MUTED)
    team = [
        ("Farjad Kareem", "CEO", "Product architecture & full-stack delivery."),
        ("Umar Azhar", "Product Manager", "Discovery through launch; MVP scope & alignment."),
        ("Irfan Ahmad", "Sr. AI Engineer", "ML / LLMs / production — owns model architecture."),
        ("Mateen Abid", "AI Engineer", "RAG pipelines, local LLM inference, integrations."),
        ("Taimoor Amir", "AI Engineer", "Whisper, NLP pipelines, FastAPI AI services."),
        ("Saad Ishaq", "Full Stack Developer", "React / TanStack, APIs, database, deployment."),
    ]
    for i, (name, role, focus) in enumerate(team):
        col = i % 3
        row = i // 3
        left = MARGIN_X + col * Inches(4.05)
        top = Inches(2.4) + row * Inches(2.05)
        card(s, left, top, Inches(3.85), Inches(1.9), name, focus, tag=role)
    add_footer(s, 13, total)

    # ── 14 CTA ───────────────────────────────────────────────
    s = blank_slide(prs)
    add_accent_rule(s, MARGIN_X, Inches(1.8), Inches(1.2))
    add_textbox(s, MARGIN_X, Inches(2.0), Inches(12), Inches(0.7), "Let's talk", size=40, bold=True, color=TEXT)
    add_textbox(
        s,
        MARGIN_X,
        Inches(2.85),
        Inches(11),
        Inches(0.8),
        "No commitment, no sales pitch — just an honest conversation about what's possible.",
        size=18,
        color=MUTED,
    )
    contacts = [
        ("Email", "info@vextoratech.com"),
        ("Phone", "+92 371 2331344"),
        ("Web", "vextoratech.com"),
        ("Location", "Lahore, Pakistan"),
    ]
    for i, (label, value) in enumerate(contacts):
        left = MARGIN_X + (i % 2) * Inches(6.1)
        top = Inches(4.0) + (i // 2) * Inches(1.15)
        add_rect(s, left, top, Inches(5.85), Inches(1.0), fill=SURFACE, line=BORDER)
        add_textbox(s, left + Inches(0.35), top + Inches(0.18), Inches(5), Inches(0.25), label.upper(), size=11, bold=True, color=PRIMARY)
        add_textbox(s, left + Inches(0.35), top + Inches(0.48), Inches(5), Inches(0.35), value, size=16, bold=True, color=TEXT)
    add_footer(s, 14, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
